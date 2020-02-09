import requests
import os
import shutil

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
from io import BytesIO
from urllib.request import urlopen
from pptx import Presentation
from pptx.util import Pt,Inches

def createSlideAPI(name, band, text):
    prs = Presentation() 
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    left = Inches(8.2)
    top = Inches(5.7)
  
    titulo = name
    subtitulo = band
    letraCompleta = text

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titulo
    subtitle.text = subtitulo
    
    letra = letraCompleta.split("\n\n")

    for linhas in letra:
        if (linhas != ''):
            letra = prs.slides.add_slide(prs.slide_layouts[2]) 
            shape = letra.shapes
            left = Inches(8.2)
            top = Inches(5.7)
            shape.add_picture("static/content/base1.png", left, top)
            body_shape = shape.placeholders[1]
            tf = body_shape.text_frame.paragraphs[0]
            tf.font.size = Pt(40)
            tf.alignment = PP_ALIGN.CENTER
            tf.font.color.rgb = RGBColor(0, 0, 0)
            tf.text += linhas
    output = BytesIO()
    prs.save(output)
    return output.getvalue()